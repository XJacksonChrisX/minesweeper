#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 建立簡報
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)  # 藍色背景
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """添加內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(236, 240, 241)  # 淺灰色背景
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(41, 128, 185)
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加雙列內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(236, 240, 241)
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(41, 128, 185)
    
    # 左列標題
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.font.size = Pt(22)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = RGBColor(52, 73, 94)
    
    # 左列內容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # 右列標題
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.font.size = Pt(22)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = RGBColor(52, 73, 94)
    
    # 右列內容
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.3), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(52, 73, 94)
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# 第 1 頁：標題頁
add_title_slide(prs, "踩地雷遊戲", "Minesweeper Game\n資工系 3B - 第 20 組")

# 第 2 頁：專案概述
add_content_slide(prs, "專案概述", [
    "• 一個使用 C++ 開發的經典踩地雷遊戲",
    "• 提供三個難度等級供玩家選擇",
    "• 透過邏輯推理和運氣揭露所有非地雷格子",
    "• 結合命令列介面與圖形化顯示",
    "• 支援 Windows、Linux、macOS 跨平台執行"
])

# 第 3 頁：成員資訊
add_content_slide(prs, "專案團隊", [
    "組長：詹書豪 (B2230041)",
    "      負責：專案規劃與整合、程式設計與測試",
    "",
    "組員：李程絛 (B2230319)",
    "",
    "組員：林元皓 (B3401214)"
])

# 第 4 頁：遊戲規則
add_content_slide(prs, "遊戲規則", [
    "1. 選擇難度等級，遊戲隨機布置地雷",
    "2. 輸入坐標揭露格子",
    "3. 格子顯示：",
    "   💣 地雷（遊戲失敗）",
    "   · 空格（周圍沒有地雷）",
    "   1-8 相鄰地雷個數",
    "4. 標記懷疑的地雷位置（🚩）",
    "5. 揭露所有非地雷格子即獲勝"
])

# 第 5 頁：難度等級
add_two_column_slide(prs, "難度等級與設定",
    "難度等級", [
        "簡單",
        "中等",
        "困難"
    ],
    "對應設定", [
        "8×8 棋盤，10 個地雷",
        "12×12 棋盤，30 個地雷",
        "16×16 棋盤，40 個地雷"
    ]
)

# 第 6 頁：系統架構 - 類別設計
add_content_slide(prs, "系統架構 - 核心類別", [
    "Cell 類別：",
    "  • 代表棋盤上的單一格子",
    "  • 屬性：是否為地雷、是否揭露、是否標記、相鄰地雷數",
    "",
    "Board 類別：",
    "  • 管理整個棋盤的狀態",
    "  • 地雷位置隨機生成、計算相鄰地雷數"
])

# 第 7 頁：系統架構 - 遊戲流程
add_content_slide(prs, "系統架構 - 遊戲流程", [
    "Game 類別：",
    "  • 控制遊戲主邏輯",
    "  • 管理遊戲狀態（進行中、勝利、失敗）",
    "  • 處理玩家輸入和遊戲進行",
    "",
    "Display 類別：",
    "  • 負責將棋盤狀態顯示給玩家"
])

# 第 8 頁：專案結構
add_two_column_slide(prs, "專案文件結構",
    "標頭檔案", [
        "board.h - 棋盤類別",
        "cell.h - 格子類別",
        "display.h - 顯示類別",
        "game.h - 遊戲類別"
    ],
    "源程式檔案", [
        "main.cpp - 主程式入口",
        "board.cpp - 棋盤實現",
        "cell.cpp - 格子實現",
        "display.cpp - 顯示實現",
        "game.cpp - 遊戲實現"
    ]
)

# 第 9 頁：編譯與執行
add_content_slide(prs, "編譯與執行", [
    "使用 Makefile 簡化編譯過程：",
    "",
    "編譯：make",
    "執行：./minesweeper（Linux/macOS）或 minesweeper.exe（Windows）",
    "清理：make clean",
    "",
    "系統要求：g++ 編譯器、C++17 標準"
])

# 第 10 頁：特色與亮點
add_content_slide(prs, "特色與亮點", [
    "✓ 跨平台支援（Windows、Linux、macOS）",
    "✓ 提供三個難度等級滿足不同玩家需求",
    "✓ 清晰的面向物件設計",
    "✓ 支援遊戲狀態保存與恢復",
    "✓ 直觀的命令列介面",
    "✓ 完整的文件與註解說明"
])

# 第 11 頁：結語
add_title_slide(prs, "感謝觀看", "踩地雷遊戲專案示範完畢")

# 保存簡報
prs.save('/workspaces/game/踩地雷遊戲_專題報告.pptx')
print("✓ PPT 檔案已成功生成：踩地雷遊戲_專題報告.pptx")
